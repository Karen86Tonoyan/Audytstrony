"""
Moduł FileGenerator - Generowanie Plików
========================================
Obsługuje generowanie:
- PDF (raporty, dokumenty)
- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)
- HTML/CSS
- Markdown
- Kod źródłowy (Python, JS, etc.)
- JSON/YAML/XML
- Obrazy z tekstem
- CSV
"""

import asyncio
import json
import csv
from pathlib import Path
from typing import Optional, List, Dict, Any, Union
from dataclasses import dataclass, field
from datetime import datetime
from enum import Enum
import io

from fpdf import FPDF
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from jinja2 import Template, Environment, FileSystemLoader
import markdown
from pygments import highlight
from pygments.lexers import get_lexer_by_name, guess_lexer
from pygments.formatters import HtmlFormatter
from PIL import Image, ImageDraw, ImageFont
from loguru import logger

from agent.config.settings import get_settings


class FileType(Enum):
    """Typy plików."""
    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    HTML = "html"
    MARKDOWN = "md"
    JSON = "json"
    YAML = "yaml"
    XML = "xml"
    CSV = "csv"
    TXT = "txt"
    PYTHON = "py"
    JAVASCRIPT = "js"
    TYPESCRIPT = "ts"
    CSS = "css"
    SQL = "sql"
    IMAGE = "png"


@dataclass
class TableData:
    """Dane tabeli."""
    headers: List[str]
    rows: List[List[Any]]
    title: Optional[str] = None


@dataclass
class ChartData:
    """Dane wykresu."""
    type: str  # bar, line, pie
    title: str
    labels: List[str]
    datasets: List[Dict[str, Any]]


@dataclass
class ReportSection:
    """Sekcja raportu."""
    title: str
    content: str
    level: int = 1  # Poziom nagłówka
    tables: List[TableData] = field(default_factory=list)
    images: List[Path] = field(default_factory=list)
    code_blocks: List[Dict[str, str]] = field(default_factory=list)


@dataclass
class Report:
    """Struktura raportu."""
    title: str
    author: str = "Ollama Agent"
    date: datetime = field(default_factory=datetime.now)
    sections: List[ReportSection] = field(default_factory=list)
    footer: Optional[str] = None
    logo: Optional[Path] = None


class PDFGenerator:
    """Generator plików PDF."""

    def __init__(self):
        self.settings = get_settings().file_generator

    def generate(self, report: Report, output_path: Path) -> Path:
        """Generuj PDF z raportu."""
        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)

        # Dodaj czcionkę z obsługą polskich znaków
        # pdf.add_font('DejaVu', '', 'DejaVuSans.ttf', uni=True)

        # Strona tytułowa
        pdf.add_page()

        if report.logo and report.logo.exists():
            pdf.image(str(report.logo), x=80, y=30, w=50)
            pdf.ln(60)

        pdf.set_font("Arial", "B", 24)
        pdf.cell(0, 20, report.title, ln=True, align="C")

        pdf.set_font("Arial", "", 14)
        pdf.cell(0, 10, f"Autor: {report.author}", ln=True, align="C")
        pdf.cell(0, 10, f"Data: {report.date.strftime('%Y-%m-%d')}", ln=True, align="C")

        # Sekcje
        for section in report.sections:
            pdf.add_page()

            # Nagłówek sekcji
            font_size = 20 - (section.level - 1) * 4
            pdf.set_font("Arial", "B", font_size)
            pdf.cell(0, 10, section.title, ln=True)
            pdf.ln(5)

            # Treść
            pdf.set_font("Arial", "", 11)
            pdf.multi_cell(0, 6, section.content)
            pdf.ln(5)

            # Tabele
            for table in section.tables:
                self._add_table(pdf, table)
                pdf.ln(10)

            # Obrazy
            for img_path in section.images:
                if img_path.exists():
                    pdf.image(str(img_path), w=150)
                    pdf.ln(10)

        # Stopka
        if report.footer:
            pdf.set_y(-30)
            pdf.set_font("Arial", "I", 8)
            pdf.cell(0, 10, report.footer, align="C")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        pdf.output(str(output_path))

        logger.info(f"Wygenerowano PDF: {output_path}")
        return output_path

    def _add_table(self, pdf: FPDF, table: TableData):
        """Dodaj tabelę do PDF."""
        if table.title:
            pdf.set_font("Arial", "B", 12)
            pdf.cell(0, 8, table.title, ln=True)

        pdf.set_font("Arial", "B", 10)
        col_width = 190 / len(table.headers)

        # Nagłówki
        for header in table.headers:
            pdf.cell(col_width, 8, str(header), border=1, align="C")
        pdf.ln()

        # Dane
        pdf.set_font("Arial", "", 9)
        for row in table.rows:
            for cell in row:
                pdf.cell(col_width, 7, str(cell)[:30], border=1)
            pdf.ln()


class WordGenerator:
    """Generator dokumentów Word."""

    def generate(self, report: Report, output_path: Path) -> Path:
        """Generuj dokument Word."""
        doc = Document()

        # Tytuł
        title = doc.add_heading(report.title, 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadane
        doc.add_paragraph(f"Autor: {report.author}")
        doc.add_paragraph(f"Data: {report.date.strftime('%Y-%m-%d')}")
        doc.add_paragraph()

        # Sekcje
        for section in report.sections:
            doc.add_heading(section.title, level=section.level)
            doc.add_paragraph(section.content)

            # Tabele
            for table_data in section.tables:
                if table_data.title:
                    doc.add_paragraph(table_data.title, style="Caption")

                table = doc.add_table(
                    rows=1 + len(table_data.rows),
                    cols=len(table_data.headers)
                )
                table.style = "Table Grid"

                # Nagłówki
                for i, header in enumerate(table_data.headers):
                    cell = table.rows[0].cells[i]
                    cell.text = str(header)
                    cell.paragraphs[0].runs[0].bold = True

                # Dane
                for row_idx, row in enumerate(table_data.rows):
                    for col_idx, value in enumerate(row):
                        table.rows[row_idx + 1].cells[col_idx].text = str(value)

            # Obrazy
            for img_path in section.images:
                if img_path.exists():
                    doc.add_picture(str(img_path), width=Inches(5))

            doc.add_paragraph()

        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))

        logger.info(f"Wygenerowano Word: {output_path}")
        return output_path


class ExcelGenerator:
    """Generator arkuszy Excel."""

    def generate(
        self,
        data: Union[TableData, List[TableData], Dict[str, TableData]],
        output_path: Path
    ) -> Path:
        """
        Generuj plik Excel.

        Args:
            data: Dane tabel (jedna tabela, lista, lub dict z nazwami arkuszy)
            output_path: Ścieżka wyjściowa
        """
        wb = Workbook()
        wb.remove(wb.active)  # Usuń domyślny arkusz

        # Normalizuj dane
        if isinstance(data, TableData):
            sheets = {"Arkusz1": data}
        elif isinstance(data, list):
            sheets = {f"Arkusz{i+1}": t for i, t in enumerate(data)}
        else:
            sheets = data

        for sheet_name, table in sheets.items():
            ws = wb.create_sheet(sheet_name)

            # Tytuł
            if table.title:
                ws.cell(row=1, column=1, value=table.title)
                ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=len(table.headers))
                ws.cell(row=1, column=1).font = Font(bold=True, size=14)
                start_row = 3
            else:
                start_row = 1

            # Nagłówki
            header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF")

            for col, header in enumerate(table.headers, 1):
                cell = ws.cell(row=start_row, column=col, value=header)
                cell.fill = header_fill
                cell.font = header_font
                cell.alignment = Alignment(horizontal="center")

            # Dane
            for row_idx, row in enumerate(table.rows, start_row + 1):
                for col_idx, value in enumerate(row, 1):
                    ws.cell(row=row_idx, column=col_idx, value=value)

            # Autofit kolumn
            for col in ws.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                ws.column_dimensions[column].width = min(max_length + 2, 50)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(output_path))

        logger.info(f"Wygenerowano Excel: {output_path}")
        return output_path


class PowerPointGenerator:
    """Generator prezentacji PowerPoint."""

    def generate(
        self,
        slides: List[Dict[str, Any]],
        output_path: Path,
        title: str = "Prezentacja"
    ) -> Path:
        """
        Generuj prezentację PowerPoint.

        Args:
            slides: Lista slajdów [{"title": "", "content": "", "image": Path, "bullets": []}]
            output_path: Ścieżka wyjściowa
            title: Tytuł prezentacji
        """
        prs = Presentation()

        # Slajd tytułowy
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = datetime.now().strftime("%Y-%m-%d")

        # Slajdy treści
        for slide_data in slides:
            if "bullets" in slide_data:
                # Layout z punktami
                bullet_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(bullet_slide_layout)
                slide.shapes.title.text = slide_data.get("title", "")

                body = slide.placeholders[1]
                tf = body.text_frame

                for i, bullet in enumerate(slide_data["bullets"]):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet

            elif "image" in slide_data:
                # Slajd z obrazem
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)

                # Tytuł
                title_shape = slide.shapes.add_textbox(
                    PptxInches(0.5), PptxInches(0.5),
                    PptxInches(9), PptxInches(1)
                )
                title_shape.text_frame.paragraphs[0].text = slide_data.get("title", "")
                title_shape.text_frame.paragraphs[0].font.size = PptxPt(32)
                title_shape.text_frame.paragraphs[0].font.bold = True

                # Obraz
                if Path(slide_data["image"]).exists():
                    slide.shapes.add_picture(
                        str(slide_data["image"]),
                        PptxInches(1), PptxInches(1.5),
                        width=PptxInches(8)
                    )

            else:
                # Zwykły slajd z treścią
                content_slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(content_slide_layout)
                slide.shapes.title.text = slide_data.get("title", "")
                slide.placeholders[1].text = slide_data.get("content", "")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

        logger.info(f"Wygenerowano PowerPoint: {output_path}")
        return output_path


class CodeGenerator:
    """Generator kodu źródłowego."""

    def generate_file(
        self,
        code: str,
        language: str,
        output_path: Path
    ) -> Path:
        """Zapisz kod do pliku."""
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(code, encoding="utf-8")

        logger.info(f"Wygenerowano kod: {output_path}")
        return output_path

    def format_code_html(
        self,
        code: str,
        language: str = None
    ) -> str:
        """Sformatuj kod jako HTML z kolorowaniem składni."""
        try:
            if language:
                lexer = get_lexer_by_name(language)
            else:
                lexer = guess_lexer(code)

            formatter = HtmlFormatter(style="monokai", linenos=True)
            return highlight(code, lexer, formatter)

        except Exception as e:
            logger.warning(f"Błąd kolorowania składni: {e}")
            return f"<pre><code>{code}</code></pre>"

    def generate_html_docs(
        self,
        code_files: Dict[str, str],
        output_path: Path,
        title: str = "Dokumentacja kodu"
    ) -> Path:
        """Generuj dokumentację HTML z kodu."""
        html_parts = [
            f"<!DOCTYPE html><html><head><title>{title}</title>",
            "<style>",
            HtmlFormatter(style="monokai").get_style_defs('.highlight'),
            "body { font-family: monospace; margin: 20px; background: #272822; color: #f8f8f2; }",
            "h1, h2 { color: #a6e22e; }",
            "</style></head><body>",
            f"<h1>{title}</h1>"
        ]

        for filename, code in code_files.items():
            lang = Path(filename).suffix.lstrip(".")
            html_parts.append(f"<h2>{filename}</h2>")
            html_parts.append(self.format_code_html(code, lang))

        html_parts.append("</body></html>")

        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text("\n".join(html_parts), encoding="utf-8")

        return output_path


class FileGeneratorModule:
    """
    Główny moduł generowania plików.
    """

    def __init__(self):
        self.settings = get_settings().file_generator
        self._pdf = PDFGenerator()
        self._word = WordGenerator()
        self._excel = ExcelGenerator()
        self._pptx = PowerPointGenerator()
        self._code = CodeGenerator()

        self._ensure_dirs()
        logger.info("FileGeneratorModule zainicjalizowany")

    def _ensure_dirs(self):
        """Upewnij się że katalogi istnieją."""
        self.settings.output_dir.mkdir(parents=True, exist_ok=True)
        self.settings.templates_dir.mkdir(parents=True, exist_ok=True)

    def _get_output_path(self, filename: str) -> Path:
        """Pobierz ścieżkę wyjściową."""
        return self.settings.output_dir / filename

    # ==================== Generatory ====================

    async def generate_pdf(
        self,
        report: Report,
        filename: str = None
    ) -> Path:
        """Generuj raport PDF."""
        if not filename:
            filename = f"raport_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output = self._get_output_path(filename)
        return self._pdf.generate(report, output)

    async def generate_word(
        self,
        report: Report,
        filename: str = None
    ) -> Path:
        """Generuj dokument Word."""
        if not filename:
            filename = f"dokument_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        output = self._get_output_path(filename)
        return self._word.generate(report, output)

    async def generate_excel(
        self,
        data: Union[TableData, List[TableData], Dict[str, TableData]],
        filename: str = None
    ) -> Path:
        """Generuj arkusz Excel."""
        if not filename:
            filename = f"arkusz_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        output = self._get_output_path(filename)
        return self._excel.generate(data, output)

    async def generate_powerpoint(
        self,
        slides: List[Dict[str, Any]],
        title: str = "Prezentacja",
        filename: str = None
    ) -> Path:
        """Generuj prezentację PowerPoint."""
        if not filename:
            filename = f"prezentacja_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        output = self._get_output_path(filename)
        return self._pptx.generate(slides, output, title)

    async def generate_html(
        self,
        content: str,
        title: str = "Strona",
        css: str = None,
        filename: str = None
    ) -> Path:
        """Generuj stronę HTML."""
        if not filename:
            filename = f"strona_{datetime.now().strftime('%Y%m%d_%H%M%S')}.html"

        html = f"""<!DOCTYPE html>
<html lang="pl">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>{css or 'body { font-family: Arial, sans-serif; margin: 40px; }'}</style>
</head>
<body>
    {content}
</body>
</html>"""

        output = self._get_output_path(filename)
        output.write_text(html, encoding="utf-8")

        logger.info(f"Wygenerowano HTML: {output}")
        return output

    async def generate_markdown(
        self,
        content: str,
        filename: str = None
    ) -> Path:
        """Generuj plik Markdown."""
        if not filename:
            filename = f"dokument_{datetime.now().strftime('%Y%m%d_%H%M%S')}.md"

        output = self._get_output_path(filename)
        output.write_text(content, encoding="utf-8")

        logger.info(f"Wygenerowano Markdown: {output}")
        return output

    async def markdown_to_html(
        self,
        md_content: str,
        title: str = "Dokument"
    ) -> str:
        """Konwertuj Markdown na HTML."""
        html_content = markdown.markdown(
            md_content,
            extensions=['tables', 'fenced_code', 'codehilite', 'toc']
        )
        return f"""<!DOCTYPE html>
<html>
<head><title>{title}</title>
<style>
body {{ font-family: Arial; max-width: 800px; margin: 0 auto; padding: 20px; }}
code {{ background: #f4f4f4; padding: 2px 5px; }}
pre {{ background: #272822; color: #f8f8f2; padding: 15px; overflow-x: auto; }}
table {{ border-collapse: collapse; width: 100%; }}
th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
th {{ background: #4CAF50; color: white; }}
</style>
</head>
<body>{html_content}</body>
</html>"""

    async def generate_json(
        self,
        data: Any,
        filename: str = None,
        indent: int = 2
    ) -> Path:
        """Generuj plik JSON."""
        if not filename:
            filename = f"data_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"

        output = self._get_output_path(filename)
        output.write_text(json.dumps(data, indent=indent, ensure_ascii=False), encoding="utf-8")

        logger.info(f"Wygenerowano JSON: {output}")
        return output

    async def generate_csv(
        self,
        data: TableData,
        filename: str = None
    ) -> Path:
        """Generuj plik CSV."""
        if not filename:
            filename = f"dane_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"

        output = self._get_output_path(filename)

        with open(output, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerow(data.headers)
            writer.writerows(data.rows)

        logger.info(f"Wygenerowano CSV: {output}")
        return output

    async def generate_code(
        self,
        code: str,
        language: str,
        filename: str = None
    ) -> Path:
        """Generuj plik z kodem."""
        extensions = {
            "python": "py", "javascript": "js", "typescript": "ts",
            "html": "html", "css": "css", "sql": "sql", "bash": "sh",
            "java": "java", "cpp": "cpp", "c": "c", "go": "go", "rust": "rs"
        }

        ext = extensions.get(language.lower(), language)
        if not filename:
            filename = f"code_{datetime.now().strftime('%Y%m%d_%H%M%S')}.{ext}"

        output = self._get_output_path(filename)
        return self._code.generate_file(code, language, output)

    async def generate_image_with_text(
        self,
        text: str,
        width: int = 800,
        height: int = 600,
        bg_color: str = "#ffffff",
        text_color: str = "#000000",
        font_size: int = 24,
        filename: str = None
    ) -> Path:
        """Generuj obraz z tekstem."""
        if not filename:
            filename = f"obraz_{datetime.now().strftime('%Y%m%d_%H%M%S')}.png"

        # Konwertuj kolory hex na RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip('#')
            return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

        img = Image.new('RGB', (width, height), hex_to_rgb(bg_color))
        draw = ImageDraw.Draw(img)

        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except:
            font = ImageFont.load_default()

        # Wyśrodkuj tekst
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        x = (width - text_width) // 2
        y = (height - text_height) // 2

        draw.text((x, y), text, fill=hex_to_rgb(text_color), font=font)

        output = self._get_output_path(filename)
        output.parent.mkdir(parents=True, exist_ok=True)
        img.save(output)

        logger.info(f"Wygenerowano obraz: {output}")
        return output

    # ==================== Szablony ====================

    async def render_template(
        self,
        template_name: str,
        context: Dict[str, Any],
        output_filename: str = None
    ) -> Path:
        """Wyrenderuj szablon Jinja2."""
        env = Environment(loader=FileSystemLoader(str(self.settings.templates_dir)))
        template = env.get_template(template_name)

        rendered = template.render(**context)

        if not output_filename:
            base = Path(template_name).stem
            ext = Path(template_name).suffix
            output_filename = f"{base}_{datetime.now().strftime('%Y%m%d_%H%M%S')}{ext}"

        output = self._get_output_path(output_filename)
        output.write_text(rendered, encoding="utf-8")

        logger.info(f"Wyrenderowano szablon: {output}")
        return output

    async def create_template(
        self,
        name: str,
        content: str
    ) -> Path:
        """Utwórz nowy szablon."""
        template_path = self.settings.templates_dir / name
        template_path.write_text(content, encoding="utf-8")

        logger.info(f"Utworzono szablon: {template_path}")
        return template_path


# ==================== Singleton ====================

_file_generator: Optional[FileGeneratorModule] = None


def get_file_generator() -> FileGeneratorModule:
    """Pobierz singleton modułu FileGenerator."""
    global _file_generator
    if _file_generator is None:
        _file_generator = FileGeneratorModule()
    return _file_generator
